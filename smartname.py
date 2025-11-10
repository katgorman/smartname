"""
SmartName - rename_files.py
A compact CLI tool to analyze mixed-media files locally (via Ollama) and propose
context-aware file name suggestions. Dry-run by default. Use --execute to apply changes.

Requirements:
  - Python 3.9+
  - pip install requests python-docx python-pptx PyMuPDF
  - Ollama running locally (https://ollama.com/) with a multimodal/text model pulled.
"""

import argparse
import base64
import json
import subprocess
import fitz  # PyMuPDF
import requests
from pathlib import Path
from docx import Document
from pptx import Presentation
import shutil

# call Ollama API
def call_ollama(model: str, prompt: str, image_data: bytes | None = None) -> str:
    url = "http://localhost:11434/api/generate"
    headers = {"Content-Type": "application/json"}
    payload = {"model": model, "prompt": prompt}
    if image_data:
        payload["images"] = [base64.b64encode(image_data).decode()]
    r = requests.post(url, headers=headers, data=json.dumps(payload), timeout=120)
    r.raise_for_status()
    text = ""
    for line in r.iter_lines():
        if line:
            chunk = json.loads(line.decode("utf-8"))
            text += chunk.get("response", "")
    return text.strip()

# extraction helperz
def extract_text(path: Path) -> str:
    ext = path.suffix.lower()
    if ext in [".txt", ".md", ".py", ".js", ".ipynb"]:
        return path.read_text(errors="ignore")[:2000]
    if ext == ".pdf":
        doc = fitz.open(path)
        text = doc[0].get_text()
        if len(text.strip()) > 50:
            return text[:2000]
        # if the pdf is scanned, render it as an image instead
        pix = doc[0].get_pixmap()
        return base64.b64encode(pix.tobytes()).decode()
    if ext == ".docx":
        doc = Document(path)
        return " ".join(p.text for p in doc.paragraphs)[:2000]
    if ext == ".pptx":
        prs = Presentation(path)
        slides = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slides.append(shape.text)
        return " ".join(slides)[:2000]
    return ""

def extract_image(path: Path) -> bytes | None:
    with open(path, "rb") as f:
        return f.read()

def extract_video_frame(path: Path) -> bytes | None:
    if not shutil.which("ffmpeg"):
        print(f"Skipping {path.name} (ffmpeg not installed)\n")
        return None
    tmp_path = path.with_suffix(".frame.jpg")
    try:
        subprocess.run(["ffmpeg", "-y", "-i", str(path), "-frames:v", "1", str(tmp_path)],
                       stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL, check=True)
        with open(tmp_path, "rb") as f:
            frame = f.read()
        tmp_path.unlink(missing_ok=True)
        return frame
    except Exception:
        print(f"Could not extract frame from {path.name}\n")
        return None

# casing engine
def apply_casing(text: str, style: str) -> str:
    words = [w.lower() for w in text.replace("_", " ").replace("-", " ").split()]
    if style == "snake_case": return "_".join(words)
    if style == "kebab-case": return "-".join(words)
    if style == "camelCase": return words[0] + "".join(w.title() for w in words[1:])
    if style == "PascalCase": return "".join(w.title() for w in words)
    if style == "lowercase": return " ".join(words)
    if style == "Title Case": return " ".join(w.title() for w in words)
    return text

# generate filenames
def generate_filename(path: Path, model: str) -> str:
    prompt = ("Analyze this file and suggest a concise, 2-5 word descriptive filename. "
              "Use underscores instead of spaces. Do not include the file extension. "
              "Respond with only the filename and nothing else.")
    ext = path.suffix.lower()
    image_data = None

    if ext in [".png", ".jpg", ".jpeg"]:
        image_data = extract_image(path)
    elif ext in [".mp4", ".mov"]:
        image_data = extract_video_frame(path)

    if image_data:
        suggestion = call_ollama(model, prompt, image_data=image_data)
    else:
        text = extract_text(path)
        suggestion = call_ollama(model, f"{prompt}\n\n{text}")
    return suggestion.strip()

# CLI entry point
def main():
    suggestions = []
    p = argparse.ArgumentParser(description="SmartName - AI-assisted file renamer")
    p.add_argument("directory", help="Target folder")
    p.add_argument("--model", default="llama3.2-vision", help="Ollama model to use")
    p.add_argument("--case", default="snake_case",
                   help="Naming style: snake_case, kebab-case, camelCase, PascalCase, lowercase, Title Case")
    p.add_argument("--execute", action="store_true", help="Apply renaming (default is dry-run)")
    args = p.parse_args()

    folder = Path(args.directory)
    if not folder.exists():
        print(f"Directory not found: {folder}")
        return

    print(f"Scanning {folder}\n")
    for f in folder.iterdir():
        if not f.is_file():
            continue
        try:
            print(f"Processing {f.name}...")
            suggestion = generate_filename(f, args.model)
            new_name = apply_casing(suggestion, args.case) + f.suffix
            suggestions.append(f,new_name)
            if args.execute:
                f.rename(f.with_name(new_name))
                print(f"Renamed to: {new_name}\n")
            else:
                print(f"Suggestion: {new_name}\n")
        except Exception as e:
            print(f"Error processing {f.name}: {e}\n")

    print("Summary of suggestions:\n")
    for original, new in suggestions:
        print(f"{original.name} -> {new}")

if __name__ == "__main__":
    main()
